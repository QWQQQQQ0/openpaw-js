"""PowerPoint presentation generator using python-pptx."""

from __future__ import annotations

from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


@dataclass
class SlideContent:
    title: str
    content: str | list[str] = field(default_factory=list)
    layout: str = "content"  # "title", "content", "two_column"


class PptGenerator:
    """Generate PowerPoint presentations from structured content."""

    def generate(
        self,
        title: str,
        slides: list[dict],
        author: str | None = None,
    ) -> bytes:
        """Generate a PowerPoint presentation.

        Args:
            title: Presentation title
            slides: List of {title: str, content: str | [str], layout?: str}
            author: Optional author metadata

        Returns:
            PPTX file as bytes
        """
        prs = Presentation()

        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Add title slide
        self._add_title_slide(prs, title, author)

        # Add content slides
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout = slide_data.get("layout", "content")

            if layout == "title":
                self._add_section_title_slide(prs, slide_title)
            elif layout == "two_column":
                self._add_two_column_slide(prs, slide_title, content)
            else:
                self._add_content_slide(prs, slide_title, content)

        # Serialize to bytes
        from io import BytesIO
        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def generate_from_markdown(
        self,
        title: str,
        markdown: str,
        author: str | None = None,
    ) -> bytes:
        """Generate presentation from Markdown content.

        Each ## heading becomes a slide title.
        Content under it becomes slide content.

        Args:
            title: Presentation title
            markdown: Markdown content
            author: Optional author

        Returns:
            PPTX file as bytes
        """
        slides = self._parse_markdown_to_slides(markdown)
        return self.generate(title, slides, author)

    def _parse_markdown_to_slides(self, markdown: str) -> list[dict]:
        """Parse Markdown into slide definitions."""
        slides = []
        current_slide = None

        for line in markdown.split("\n"):
            line = line.rstrip()

            # H2 = new slide
            if line.startswith("## "):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "title": line[3:].strip(),
                    "content": [],
                }
            elif line.startswith("### "):
                # H3 = section title slide
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "title": line[4:].strip(),
                    "content": [],
                    "layout": "title",
                }
            elif current_slide is not None:
                if line.strip():
                    current_slide["content"].append(line.strip())

        if current_slide:
            slides.append(current_slide)

        # Convert content lists to strings
        for slide in slides:
            if isinstance(slide["content"], list):
                slide["content"] = "\n".join(slide["content"])

        return slides

    def _add_title_slide(self, prs: Presentation, title: str, author: str | None) -> None:
        """Add the title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(11)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Add author
        if author:
            top = Inches(4.2)
            height = Inches(1)
            author_box = slide.shapes.add_textbox(left, top, width, height)
            author_frame = author_box.text_frame
            author_frame.text = f"作者: {author}"
            author_para = author_frame.paragraphs[0]
            author_para.font.size = Pt(20)
            author_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            author_para.alignment = PP_ALIGN.CENTER

    def _add_section_title_slide(self, prs: Presentation, title: str) -> None:
        """Add a section title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        left = Inches(1)
        top = Inches(3)
        width = Inches(11)
        height = Inches(1.5)
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.text = title
        para = frame.paragraphs[0]
        para.font.size = Pt(36)
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, title: str, content: str | list[str]) -> None:
        """Add a content slide with title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Content
        top = Inches(1.5)
        height = Inches(5.5)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        if isinstance(content, str):
            lines = content.split("\n")
        else:
            lines = content

        for i, line in enumerate(lines):
            if i == 0:
                content_frame.text = line
                content_frame.paragraphs[0].font.size = Pt(18)
            else:
                para = content_frame.add_paragraph()
                para.text = line
                para.font.size = Pt(18)
                para.space_before = Pt(8)

    def _add_two_column_slide(self, prs: Presentation, title: str, content: str | list[str]) -> None:
        """Add a two-column slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Split content into two columns
        if isinstance(content, str):
            lines = content.split("\n")
        else:
            lines = content

        mid = len(lines) // 2
        col1_lines = lines[:mid]
        col2_lines = lines[mid:]

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, line in enumerate(col1_lines):
            if i == 0:
                left_frame.text = line
                left_frame.paragraphs[0].font.size = Pt(16)
            else:
                para = left_frame.add_paragraph()
                para.text = line
                para.font.size = Pt(16)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(5.5), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, line in enumerate(col2_lines):
            if i == 0:
                right_frame.text = line
                right_frame.paragraphs[0].font.size = Pt(16)
            else:
                para = right_frame.add_paragraph()
                para.text = line
                para.font.size = Pt(16)


# Standalone usage
if __name__ == "__main__":
    gen = PptGenerator()
    md = """
## 项目背景

- 当前市场环境分析
- 用户需求调研结果
- 竞品对比

## 核心功能

- 功能一：智能文档生成
- 功能二：数据可视化
- 功能三：自动化流程

## 技术架构

- 前端：React + TypeScript
- 后端：Python + FastAPI
- 数据库：PostgreSQL

## 未来规划

- Q1：完成核心功能开发
- Q2：用户测试与反馈
- Q3：正式发布
"""
    pptx_bytes = gen.generate_from_markdown("项目介绍", md, author="OpenPaw")
    with open("test_output.pptx", "wb") as f:
        f.write(pptx_bytes)
    print("Generated test_output.pptx")
